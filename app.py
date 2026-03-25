import streamlit as st
import chromadb
from chromadb.utils.embedding_functions import DefaultEmbeddingFunction
from pathlib import Path
import docx
import pptx

# --- Configuração ---
DB_PATH = str(Path.home() / "chromadb-demo" / "chroma_db")
COLLECTION_NAME = "tjmg_docs"
PASTA_DOCS = str(Path.home() / "chromadb-demo" / "docs")

st.set_page_config(page_title="Busca Semântica - TJMG", page_icon="🔍", layout="wide")
st.title("🔍 Busca Semântica de Documentos — TJMG")

# --- Cache de cliente e função de embedding ---
@st.cache_resource
def carregar_recursos():
    ef = DefaultEmbeddingFunction()
    client = chromadb.PersistentClient(path=DB_PATH)
    collection = client.get_or_create_collection(COLLECTION_NAME, embedding_function=ef)
    return client, collection, ef

client, collection, ef = carregar_recursos()

# --- Funções de extração de texto ---
def extrair_texto_docx(caminho):
    doc = docx.Document(caminho)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extrair_texto_pptx(caminho):
    prs = pptx.Presentation(caminho)
    textos = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                textos.append(shape.text)
    return "\n".join(textos)

def extrair_texto(caminho):
    ext = Path(caminho).suffix.lower()
    if ext == ".docx":
        return extrair_texto_docx(caminho)
    elif ext == ".pptx":
        return extrair_texto_pptx(caminho)
    elif ext in [".txt", ".md"]:
        return Path(caminho).read_text(encoding="utf-8", errors="ignore")
    return None

def chunk_texto(texto, tamanho=500, sobreposicao=50):
    palavras = texto.split()
    chunks = []
    i = 0
    while i < len(palavras):
        chunk = " ".join(palavras[i:i + tamanho])
        chunks.append(chunk)
        i += tamanho - sobreposicao
    return [c for c in chunks if len(c.strip()) > 30]

# --- Sidebar ---
with st.sidebar:
    st.header("Configuração")
    pasta_input = st.text_input("Pasta de documentos", value=PASTA_DOCS)
    st.caption("Formatos suportados: .docx, .pptx, .txt, .md")

    st.divider()
    total = collection.count()
    st.metric("Trechos indexados", total)

    if st.button("Indexar / Atualizar", type="primary", use_container_width=True):
        pasta = Path(pasta_input)
        if not pasta.exists():
            pasta.mkdir(parents=True, exist_ok=True)
            st.warning(f"Pasta criada: {pasta}\nAdicione documentos e clique novamente.")
        else:
            arquivos = (list(pasta.rglob("*.docx")) + list(pasta.rglob("*.pptx")) +
                        list(pasta.rglob("*.txt")) + list(pasta.rglob("*.md")))

            if not arquivos:
                st.warning("Nenhum arquivo encontrado na pasta.")
            else:
                barra = st.progress(0, text="Indexando...")
                adicionados = 0
                erros = 0

                for i, arq in enumerate(arquivos):
                    try:
                        texto = extrair_texto(str(arq))
                        if not texto or len(texto.strip()) < 20:
                            continue
                        chunks = chunk_texto(texto)
                        for j, chunk in enumerate(chunks):
                            doc_id = f"{arq.stem}_{j}"
                            collection.upsert(
                                ids=[doc_id],
                                documents=[chunk],
                                metadatas=[{"arquivo": arq.name, "chunk": j}]
                            )
                        adicionados += 1
                    except Exception as e:
                        erros += 1
                        st.error(f"Erro em {arq.name}: {e}")
                    barra.progress((i + 1) / len(arquivos), text=f"Processando {arq.name}...")

                barra.empty()
                st.success(f"{adicionados} arquivo(s) indexado(s). Erros: {erros}")
                st.rerun()

    if st.button("Limpar índice", use_container_width=True):
        client.delete_collection(COLLECTION_NAME)
        st.cache_resource.clear()
        st.success("Índice limpo.")
        st.rerun()

# --- Área de busca ---
st.subheader("Buscar documentos")
consulta = st.text_input("Digite sua busca", placeholder="Ex: nota fiscal de material de limpeza")

col1, col2 = st.columns([1, 4])
with col1:
    n_resultados = st.number_input("Resultados", min_value=1, max_value=20, value=5)

if consulta:
    if collection.count() == 0:
        st.warning("Nenhum documento indexado. Use o botão 'Indexar / Atualizar' na barra lateral.")
    else:
        with st.spinner("Buscando..."):
            resultados = collection.query(
                query_texts=[consulta],
                n_results=min(n_resultados, collection.count()),
                include=["documents", "metadatas", "distances"]
            )

        docs = resultados["documents"][0]
        metas = resultados["metadatas"][0]
        distancias = resultados["distances"][0]

        if not docs:
            st.info("Nenhum resultado encontrado.")
        else:
            st.write(f"**{len(docs)} resultado(s) encontrado(s):**")
            for i, (doc, meta, dist) in enumerate(zip(docs, metas, distancias)):
                similaridade = round((1 - dist) * 100, 1)
                with st.expander(
                    f"#{i+1} — {meta.get('arquivo', '?')}  (similaridade: {similaridade}%)"
                ):
                    st.write(doc)
                    st.caption(f"Arquivo: {meta.get('arquivo')} | Chunk: {meta.get('chunk')}")
